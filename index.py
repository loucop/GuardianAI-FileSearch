import os
import json
import sqlite3
import platform
import hashlib
import argparse
from concurrent.futures import ThreadPoolExecutor, ProcessPoolExecutor, as_completed
from tqdm import tqdm

# Importações de leitores de arquivo
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pdf2image import convert_from_path
import pytesseract
from pptx import Presentation

# Embedding local
from sentence_transformers import SentenceTransformer

# ChromaDB
import chromadb

# ==============================================================================  
# CONFIGURAÇÃO  
# ==============================================================================  
# ATENÇÃO: Configure estes caminhos antes de executar o script.  

# (Opcional, apenas para Windows) Caminho para o executável do Tesseract OCR.  
TESSERACT_CMD_PATH = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# Diretório onde o banco de dados do ChromaDB será armazenado.
CHROMA_PERSIST_DIRECTORY = "./chroma_db"

# Diretório do banco de dados de fallback (SQLite).
SQLITE_DB_PATH = "embeddings_fallback.db"

# ==============================================================================  

# Configura o Tesseract se o caminho foi fornecido
if TESSERACT_CMD_PATH and os.path.exists(TESSERACT_CMD_PATH):
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD_PATH

# Constantes
IMAGE_EXTS = [".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif"]
MAX_THREAD_WORKERS = 4
MAX_PROCESS_WORKERS = os.cpu_count() or 2  # Usa os cores disponíveis

# --------- Modelo de Embedding (local) ----------
print("Carregando modelo de embeddings local (all-MiniLM-L6-v2)...")
sbert = SentenceTransformer("all-MiniLM-L6-v2")

def embed_texts(texts):
    """Gera embeddings para uma lista de textos."""
    embs = sbert.encode(texts, show_progress_bar=False, convert_to_tensor=False)
    # Garante que a saída seja uma lista de listas de floats
    return [e.tolist() if hasattr(e, "tolist") else list(map(float, e)) for e in embs]

# --------- Funções do Banco de Dados (ChromaDB e Fallback SQLite) ----------

def make_id(text: str) -> str:
    """Cria um ID determinístico a partir de um texto (ex: caminho do arquivo)."""
    return hashlib.sha1(text.encode("utf-8")).hexdigest()

class DatabaseManager:
    """
    Gerencia o armazenamento em ChromaDB com fallback para SQLite.
    Agora lida com múltiplas coleções.
    """
    def __init__(self, use_chroma=True):
        self.using_chroma = False
        if use_chroma:
            try:
                print(f"Inicializando ChromaDB em '{CHROMA_PERSIST_DIRECTORY}'...")
                self.chroma_client = chromadb.PersistentClient(path=CHROMA_PERSIST_DIRECTORY)
                self.collections = {
                    "files_content": self.chroma_client.get_or_create_collection("files_content"),
                    "files_name": self.chroma_client.get_or_create_collection("files_name"),
                    "folders": self.chroma_client.get_or_create_collection("folders"),
                }
                self.using_chroma = True
                print("ChromaDB inicializado com sucesso.")
            except Exception as e:
                print(f"AVISO: Falha ao inicializar ChromaDB: {e}")
                self.using_chroma = False

        if not self.using_chroma:
            print(f"Usando fallback para SQLite em '{SQLITE_DB_PATH}'.")
            self.sqlite_init()

    def sqlite_init(self):
        """Inicializa as tabelas no banco de dados SQLite."""
        with sqlite3.connect(SQLITE_DB_PATH) as conn:
            cur = conn.cursor()
            cur.execute("CREATE TABLE IF NOT EXISTS files_content (id TEXT PRIMARY KEY, path TEXT, content TEXT, embedding TEXT)")
            cur.execute("CREATE TABLE IF NOT EXISTS files_name (id TEXT PRIMARY KEY, path TEXT, name TEXT, embedding TEXT)")
            cur.execute("CREATE TABLE IF NOT EXISTS folders (id TEXT PRIMARY KEY, path TEXT, embedding TEXT)")
            conn.commit()

    def get_existing_ids(self, collection_name: str) -> set:
        """Obtém os IDs existentes de uma coleção/tabela."""
        if self.using_chroma:
            try:
                return set(self.collections[collection_name].get(include=[])["ids"])
            except Exception as e:
                print(f"Erro ao buscar IDs do Chroma para '{collection_name}': {e}")
                return set()
        else:  # SQLite
            with sqlite3.connect(SQLITE_DB_PATH) as conn:
                cur = conn.cursor()
                cur.execute(f"SELECT id FROM {collection_name}")
                return {row[0] for row in cur.fetchall()}

    def store_batch(self, collection_name: str, data: list):
        """Armazena um lote de dados na coleção/tabela apropriada."""
        if not data:
            return

        ids = [item['id'] for item in data]
        embeddings = [item['embedding'] for item in data]
        documents = [item['document'] for item in data]
        metadatas = [item['metadata'] for item in data]

        if self.using_chroma:
            try:
                self.collections[collection_name].add(
                    ids=ids,
                    embeddings=embeddings,
                    documents=documents,
                    metadatas=metadatas
                )
            except Exception as e:
                print(f"Erro ao adicionar lote no Chroma ('{collection_name}'), tentando fallback para SQLite: {e}")
                if collection_name == "files_content":
                     self._sqlite_store_content(ids, documents, embeddings, metadatas)
        elif collection_name == "files_content":
            self._sqlite_store_content(ids, documents, embeddings, metadatas)

    def _sqlite_store_content(self, ids, documents, embeddings, metadatas):
        """Lógica específica para salvar conteúdo de arquivos no SQLite."""
        with sqlite3.connect(SQLITE_DB_PATH) as conn:
            cur = conn.cursor()
            rows = []
            for id_val, doc, emb, meta in zip(ids, documents, embeddings, metadatas):
                rows.append((id_val, meta.get('path', ''), doc, json.dumps(emb)))
            cur.executemany("INSERT OR REPLACE INTO files_content (id, path, content, embedding) VALUES (?, ?, ?, ?)", rows)
            conn.commit()

# --------- Funções de Leitura de Arquivo e OCR ----------
def read_file_content(file_path, use_ocr=True):
    """Lê o conteúdo de um arquivo, com opção de usar OCR para PDFs e imagens."""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".txt":
            try:
                with open(file_path, "r", encoding="utf-8") as f: return f.read()
            except UnicodeDecodeError:
                with open(file_path, "r", encoding="latin-1") as f: return f.read()
        elif ext == ".csv":
            return pd.read_csv(file_path).to_string()
        elif ext in [".xls", ".xlsx", ".xlsm", ".xlsb"]:
            df = pd.read_excel(file_path, sheet_name=None)
            return "\n".join([f"--- {name} ---\n{sheet.to_string()}" for name, sheet in df.items()])
        elif ext == ".docx":
            return "\n".join([p.text for p in Document(file_path).paragraphs])
        elif ext == ".pptx":
            return "\n".join(shape.text for slide in Presentation(file_path).slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext == ".pdf":
            text = ""
            try:
                reader = PdfReader(file_path)
                for page in reader.pages: text += page.extract_text() or ""
            except Exception: pass
            if use_ocr and (not text or not text.strip()):
                images = convert_from_path(file_path)
                for img in images: text += pytesseract.image_to_string(img) or ""
            return text
        elif ext in IMAGE_EXTS and use_ocr:
            from PIL import Image
            return pytesseract.image_to_string(Image.open(file_path))
        return ""
    except Exception:
        return ""

def process_file_worker(path):
    """Worker que lê o conteúdo de um arquivo e retorna o caminho e o conteúdo."""
    content = read_file_content(path, use_ocr=True)
    return path, content

# --------- Indexador Principal ----------
def index_folder(root_folder, batch_size=32, resume=True):
    """
    Indexa todos os arquivos em um diretório, criando embeddings para conteúdo,
    nomes de arquivos e nomes de pastas.
    """
    db_manager = DatabaseManager(use_chroma=True)

    print("Procurando por arquivos...")
    all_files = [os.path.join(root, fn) for root, _, files in os.walk(root_folder) for fn in files]
    print(f"Encontrados {len(all_files)} arquivos no total.")

    # Filtrar arquivos já processados (se `resume` for True)
    if resume:
        print("Verificando arquivos já indexados...")
        existing_content_ids = db_manager.get_existing_ids("files_content")
        files_to_process = [f for f in all_files if make_id(f) not in existing_content_ids]
        print(f"{len(files_to_process)} novos arquivos para processar.")
    else:
        files_to_process = all_files

    if not files_to_process:
        print("Nenhum arquivo novo para indexar.")
        return

    # 1. Processar CONTEÚDO dos arquivos
    content_batch = []
    with ProcessPoolExecutor(max_workers=MAX_PROCESS_WORKERS) as executor:
        futures = [executor.submit(process_file_worker, f) for f in files_to_process]
        for future in tqdm(as_completed(futures), total=len(futures), desc="Lendo conteúdo dos arquivos"):
            try:
                path, content = future.result()
                if content and content.strip():
                    content_batch.append({'path': path, 'content': content})
                    if len(content_batch) >= batch_size:
                        contents = [item['content'] for item in content_batch]
                        embeddings = embed_texts(contents)
                        data_to_store = [{
                            'id': make_id(item['path']),
                            'embedding': emb,
                            'document': item['content'],
                            'metadata': {'path': item['path']}
                        } for item, emb in zip(content_batch, embeddings)]
                        db_manager.store_batch("files_content", data_to_store)
                        content_batch = []
            except Exception as e:
                print(f"Erro no worker de arquivo: {e}")

    if content_batch:
        contents = [item['content'] for item in content_batch]
        embeddings = embed_texts(contents)
        data_to_store = [{
            'id': make_id(item['path']),
            'embedding': emb,
            'document': item['content'],
            'metadata': {'path': item['path']}
        } for item, emb in zip(content_batch, embeddings)]
        db_manager.store_batch("files_content", data_to_store)

    # 2. Processar NOMES de arquivos e PASTAS (para todos os arquivos, com batches)
    print("Indexando nomes de arquivos e pastas...")

    # Nomes de arquivos
    file_names = {os.path.basename(f): f for f in all_files}
    unique_file_names = list(file_names.keys())
    BATCH_SIZE_NAME = 5000
    for i in range(0, len(unique_file_names), BATCH_SIZE_NAME):
        batch_names = unique_file_names[i:i+BATCH_SIZE_NAME]
        batch_embeddings = embed_texts(batch_names)
        data_to_store_names = [{
            'id': make_id(name),
            'embedding': emb,
            'document': name,
            'metadata': {'path': file_names[name]}
        } for name, emb in zip(batch_names, batch_embeddings)]
        db_manager.store_batch("files_name", data_to_store_names)

    print(f"{len(unique_file_names)} nomes de arquivos únicos indexados.")

    # Nomes de pastas
    folder_paths = {os.path.dirname(f) for f in all_files}
    unique_folder_paths = list(folder_paths)
    BATCH_SIZE_FOLDER = 5000
    for i in range(0, len(unique_folder_paths), BATCH_SIZE_FOLDER):
        batch_folders = unique_folder_paths[i:i+BATCH_SIZE_FOLDER]
        batch_embeddings = embed_texts(batch_folders)
        data_to_store_folders = [{
            'id': make_id(path),
            'embedding': emb,
            'document': path,
            'metadata': {'path': path}
        } for path, emb in zip(batch_folders, batch_embeddings)]
        db_manager.store_batch("folders", data_to_store_folders)

    print(f"{len(unique_folder_paths)} pastas únicas indexadas.")

    print("\nIndexação finalizada!")

def main():
    parser = argparse.ArgumentParser(description="Indexador de arquivos para busca semântica.")
    parser.add_argument("root_folder", type=str, help="O caminho para a pasta raiz que você deseja indexar.")
    parser.add_argument("--no-resume", action="store_true", help="Força a reindexação de todos os arquivos, ignorando o cache.")
    parser.add_argument("--batch-size", type=int, default=32, help="Número de arquivos para processar em cada lote.")
    
    args = parser.parse_args()
    
    if not os.path.isdir(args.root_folder):
        print(f"Erro: O diretório '{args.root_folder}' não foi encontrado.")
        return

    index_folder(args.root_folder, batch_size=args.batch_size, resume=not args.no_resume)

if __name__ == "__main__":
    main()
